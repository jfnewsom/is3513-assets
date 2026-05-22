"""
build_lib.py - Shared deck assembly library for IS3513 slide decks.

Per-deck build scripts (e.g., slides/m1_c1/build_m1_c1.py) import DeckBuilder
and the layout-specific add_* methods, then make a sequence of calls with the
deck's content inline as string literals. DeckBuilder handles the template
unpack, slide duplication, content population, PNG render-and-insert, speaker
notes, and final packaging.

PATTERN
-------
    from build_lib import DeckBuilder, format_concept_notes, format_title_notes

    deck = DeckBuilder(png_out='/home/claude/m1_c1_pngs',
                       work='/tmp/build_m1_c1')

    deck.add_title_slide(
        course_id='IS3513',
        course_name='Information Assurance and Security',
        subtitle='Module 1, Chapter 1: Security Trends',
        attribution='PROF. JOHN NEWSOM   ·   SUMMER 2026   ·   SECTION 0XX',
        notes=format_title_notes('M1-C1', 'Security Trends', ...),
    )

    deck.add_concept_slide(
        kicker='topic',
        title='What This Chapter Is About',
        subhead='The threat landscape in one paragraph',
        section_kicker='THE BIG IDEA',
        card_heading='Threats keep changing, fundamentals do not.',
        lead='One or two sentences setting up what this slide is about.',
        bullets=[
            'Threat actors evolve, but motivations stay the same.',
            'Defense in depth still applies.',
            'CIA triad is still the scoring rubric.',
            'New tooling changes the speed, not the shape.',
        ],
        notes=format_concept_notes(video_script='...', think_about=[...]),
    )

    # ... more slides ...

    deck.save('/home/claude/M1-C1.pptx')


LAYOUT TO TEMPLATE SOURCE MAPPING (v6 template, 12 source slides)
-----------------------------------------------------------------
    title       -> slide 1   (course branding + welcome)
    concept     -> slide 2   (single-card with kicker, heading, lead, 4 bullets)
    two_col     -> slide 3   (two-column comparison, 4 bullets per side)
    intro       -> slide 4   (portrait + facts + pull-quote)
    team        -> slide 5   (5-person mentor grid)
    stats       -> slide 6   (3-stat bar + supporting bullets)
    structure   -> slide 7   (three-card row, e.g. Unit 1/2/3)
    warning     -> slide 8   (hard-truth callout with supporting points)
    support     -> slide 9   (contact channels grid)
    module      -> slide 10  (module overview with 3-lab grid)
    guide       -> slide 11  (three stacked callouts: Do This / Heads Up / FYI)
    policy      -> slide 12  (four numbered rules)

When the template gains new layouts, add the slide source index to LAYOUT
and write a new add_*_slide method.
"""
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

# ----- Paths -----
LIB_DIR = Path(__file__).resolve().parent
REPO_ROOT = LIB_DIR.parent
TOOLS_DIR = REPO_ROOT / 'tools'
DEFAULT_TEMPLATE = REPO_ROOT / 'ppt' / 'IS3513_Template_v7.pptx'

# Make the renderer importable
sys.path.insert(0, str(TOOLS_DIR))
from render_slides import render_code, render_output  # noqa: E402

# pptx skill scripts
PPTX_SCRIPTS = Path('/mnt/skills/public/pptx/scripts')
UNPACK = PPTX_SCRIPTS / 'office' / 'unpack.py'
PACK = PPTX_SCRIPTS / 'office' / 'pack.py'
CLEAN = PPTX_SCRIPTS / 'clean.py'
ADD_SLIDE = PPTX_SCRIPTS / 'add_slide.py'

# ----- Template layout source indices (1-based template slide numbers) -----
LAYOUT = {
    'title':     1,
    'concept':   2,
    'two_col':   3,
    'intro':     4,
    'team':      5,
    'stats':     6,
    'structure': 7,
    'warning':   8,
    'support':   9,
    'module':    10,
    'guide':     11,
    'policy':    12,
}

# ----- Placeholder text in each template slide (used for surgical replacement) -----
# Strings here MUST match the literal text in the v6 template's XML.
# If the template changes, regenerate this dict by running:
#     python3 -c "from pptx import Presentation; ... " against the template.
#
# Note: '&' in the template XML is '&amp;'. The replace helper handles XML escaping
# of new content automatically, but old strings must be supplied as the user-visible
# text (e.g., 'Info Bar & Stats Layout', not 'Info Bar &amp; Stats Layout').
# Smart quotes (\u201C \u201D) are passed through.
PH = {
    'title': {
        'course_id':    'IS3513',
        'course_name':  'Information Assurance and Security',
        'subtitle':     'Welcome: NEXUS Security Operations',
        'attribution':  'PROF. JOHN NEWSOM   \u00b7   SUMMER 2026   \u00b7   SECTION 0XX',
    },
    'concept': {
        'kicker':         'topic',
        'title':          'Slide Title Goes Here',
        'subhead':        'Subhead or context line in the yellow band',
        'section_kicker': 'SECTION KICKER (H3 / UPPERCASE / ACCENT)',
        'card_heading':   'Card heading: short and direct',
        'lead':           "Lead paragraph. One or two sentences setting up what this slide is about. Body text uses Roboto in the muted gray that matches the site's body color on dark cards.",
        'bullets': [
            'Bullet point one: concise, action-oriented phrasing',
            'Bullet point two: student-facing language, conversational where it fits',
            'Bullet point three: keep each bullet to about one line at this size',
            'Bullet point four: break into a second slide when you hit five or six items',
        ],
    },
    'two_col': {
        'kicker':         'topic',
        'title':          'Two-Column Comparison',
        'subhead':        'Side-by-side lists for parallel content',
        'left_kicker':    'LEFT COLUMN KICKER',
        'left_heading':   'What this is for',
        'left_bullets': [
            'First parallel item left',
            'Second parallel item left',
            'Third parallel item left',
            'Fourth parallel item left',
        ],
        'right_kicker':   'RIGHT COLUMN KICKER',
        'right_heading':  "What this isn't",
        'right_bullets': [
            'First parallel item right',
            'Second parallel item right',
            'Third parallel item right',
            'Fourth parallel item right',
        ],
    },
    'intro': {
        'kicker':       'intro',
        'title':        'Portrait + Facts Layout',
        'subhead':      'For mentor profiles or instructor intro',
        'photo_label':  '3.0 \u00d7 3.0 IN',
        'role_kicker':  'ROLE OR TAG',
        'name':         'Person Name',
        'tagline':      'One-line tagline in italic, sets the voice.',
        'facts': [
            'Title / position / affiliation',
            'Background or experience point',
            'Course role: what they teach or coach',
            'Reach: how students contact them',
        ],
        # The pull-quote contains entity-encoded smart quotes in v6's XML.
        # We match it via _replace_text_raw (skips the &-escape).
        'pullquote_raw': "&#x201C;A short pull-quote that captures this person's perspective.&#x201D;",
    },
    'team': {
        'kicker':   'team',
        'title':    'Meet Your Mentors',
        'subhead':  "Five voices you'll hear all semester",
        # Each member is (badge, name_ph, role_ph). The badges (#1..#5) are
        # unique strings and serve as the surgical anchors.
        'members': [
            ('#1', 'First Lastname', 'Role / Title'),
            ('#2', 'First Lastname', 'Role / Title'),
            ('#3', 'First Lastname', 'Role / Title'),
            ('#4', 'First Lastname', 'Role / Title'),
            ('#5', 'First Lastname', 'Role / Title'),
        ],
        'framing':  'Short framing line that sets up why these people matter.',
    },
    'stats': {
        'kicker':         'stats',
        'title':          'Info Bar & Stats Layout',
        'subhead':        'Quick numeric summary at the top',
        'stats': [
            ('MODULES',             '5'),
            ('FOUNDATION LABS',     '9'),
            ('ENGAGEMENT PACKETS',  '5'),
        ],
        'lead':           "Use this layout when the slide's job is to anchor a quantitative claim: course structure, grading weights, time commitment.",
        'bullets': [
            'Each stat sits in its own bordered tile with centered content',
            'Values use Roboto Slab in the slide accent color',
            'Keep to 3 or 4 stats: more crowds the row',
        ],
    },
    'structure': {
        'kicker':   'structure',
        'title':    'Three-Card Row',
        'subhead':  'For Unit 1 / Unit 2 / Unit 3 or any 3-way split',
        'lead':     'Each module breaks into three units. Each unit has one Foundation Lab; the Engagement Packet sits inside Unit 3.',
        'cards': [
            ('UNIT 1', 'Training',
             'Build the foundation. Concepts, environment setup, and first hands-on work.'),
            ('UNIT 2', 'Training+',
             'Apply the foundation. Deeper exercises, more autonomy, fewer guardrails.'),
            ('UNIT 3', 'Engagement',
             "Client work. Combine the unit's skills into the Engagement Packet deliverable."),
        ],
    },
    'warning': {
        'kicker':       'warning',
        'title':        'Warning / Hard-Truth Slide',
        'subhead':      'Use sparingly: when weight matters',
        'banner':       'READ THIS TWICE',
        'banner_line':  'A statement that needs to land',
        'rule_kicker':  'THE RULE',
        'rule_oneliner':'Stated in one sentence, no softening.',
        'rule_body':    "Explanation of the rule, the rationale, and what happens if it's ignored. Body text stays calm even when the topic is sharp.",
        'bullets': [
            'Supporting point one: concrete consequence or example',
            'Supporting point two: what the safety net looks like',
            'Supporting point three: what counts as the line',
        ],
    },
    'support': {
        'kicker':         'support',
        'title':          'Contact & Support Channels',
        'subhead':        'How to reach me when you need help',
        'philo_kicker':   'DISCORD FIRST, EMAIL FOR ANYTHING PERSONAL',
        'philo_heading':  'How we talk to each other',
        'philo_body':     'Discord for general questions and peer help. Email for anything personal or FERPA-sensitive. Calendly for 1-on-1 time. Tuesday 6 PM Central drop-in Zoom.',
        'channels': [
            ('DISCORD',   'Post in #course-questions'),
            ('EMAIL',     'john.newsom@utsa.edu'),
            ('CALENDLY',  '1:1 by appointment'),
            ('ZOOM',      'Tuesdays 6 PM Central'),
        ],
    },
    'module': {
        'number':       '0',
        'title':        'Module title in the yellow band',
        'labs_label':   'LABS',
        'labs_value':   '3 Labs',
        'chap_label':   'CHAPTERS',
        'chap_value':   'Ch X & Y',
        'client_label': 'CLIENT',
        'client_value': 'Client Name',
        'overview':     'Module overview: one or two sentences setting up what this module covers and how it connects to the client engagement at the end.',
        'labs': [
            ('LAB X.1', 'Foundation lab title goes here', 'Internal Training'),
            ('LAB X.2', 'Foundation lab title goes here', 'Internal Training'),
            ('LAB X.3', 'Engagement Packet title here',   'Client Billable'),
        ],
        # Slide 10 has a 'CLIENT\nLOGO' image placeholder; leave alone.
    },
    'guide': {
        'kicker':   'guide',
        'title':    'Stacked Callouts',
        'subhead':  'When the content is a sequence of weighted points',
        'callouts': [
            ('DO THIS',         'Practice as you go',     'First action item. Specific, time-bounded, achievable in the next 24 hours.'),
            ('HEADS UP',        'One thing to watch',     "A caution that's not a hard rule. Something that trips students up if they don't notice."),
            ('USEFUL TO KNOW',  'Why this matters',       "Background context that helps the rest stick. Goes last because it's framing, not action."),
        ],
    },
    'policy': {
        'kicker':         'policy',
        'title':          'Rules & Acknowledgment',
        'subhead':        'Numbered rules with an ack box',
        'rules_kicker':   'THE FOUR RULES',
        # Each rule is (number, title_placeholder, body_placeholder). The numbers
        # (1..4) are unique short strings; the title and body strings are repeated
        # across rules, so we use the per-rule number as an anchor when replacing.
        # Because the title/body strings are identical, we have to do four
        # sequential replacements with re.subn count=1 in the order rules appear
        # in the XML.
        'rules_title_repeat': 'Rule title here',
        'rules_body_repeat':  'One-sentence explanation of what this rule means in practice.',
    },
}

# ----- PNG placement (when a demo/output PNG overlay is added) -----
PNG_X = Inches(0.75)
PNG_Y = Inches(2.20)
PNG_W = Inches(11.83)


# =============================================================================
# Helpers
# =============================================================================

def _xml_escape(s):
    return s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')


def _replace_text(xml, old_plain, new_plain):
    """Replace plaintext inside an <a:t> tag. Old/new are user-visible strings;
    they are XML-escaped before/after the match."""
    old_escaped = _xml_escape(old_plain)
    new_escaped = _xml_escape(new_plain)
    pattern = rf'(<a:t(?:\s[^>]*)?>){re.escape(old_escaped)}(</a:t>)'
    new_xml, count = re.subn(pattern, rf'\g<1>{new_escaped}\g<2>', xml)
    if count == 0:
        raise ValueError(f'Placeholder not found in slide XML: {old_plain!r}')
    return new_xml


def _replace_text_raw(xml, old_raw, new_plain):
    """Replace text whose old form already contains XML entities (e.g. &#x201C;).
    Only the new content is XML-escaped."""
    new_escaped = _xml_escape(new_plain)
    pattern = rf'(<a:t(?:\s[^>]*)?>){re.escape(old_raw)}(</a:t>)'
    new_xml, count = re.subn(pattern, rf'\g<1>{new_escaped}\g<2>', xml)
    if count == 0:
        raise ValueError(f'Raw placeholder not found in slide XML: {old_raw!r}')
    return new_xml


def _replace_first(xml, old_plain, new_plain):
    """Replace only the first occurrence of old_plain. Used for repeated
    placeholders that must be replaced in document order (e.g. policy rules)."""
    old_escaped = _xml_escape(old_plain)
    new_escaped = _xml_escape(new_plain)
    pattern = rf'(<a:t(?:\s[^>]*)?>){re.escape(old_escaped)}(</a:t>)'
    new_xml, count = re.subn(pattern, rf'\g<1>{new_escaped}\g<2>', xml, count=1)
    if count == 0:
        raise ValueError(f'Placeholder not found in slide XML: {old_plain!r}')
    return new_xml


def _inject_extra_bullet(xml, text):
    """Inject a 5th bullet paragraph into the concept-slide bullet text-body.

    The concept slide stores all 4 bullets as <a:p> paragraphs inside a single
    <p:txBody>. We locate the last (4th) bullet paragraph and append a clone of
    it immediately after, with the placeholder text swapped for the supplied
    text. The new paragraph inherits the formatting of bullet 4.

    Used by add_concept_slide(extra_bullet=...). Standard concept slides have
    4 bullets; this is an explicit exception for slides that genuinely need 5.
    """
    bullet4_ph = PH['concept']['bullets'][3]
    pattern = re.compile(
        r'(<a:p>(?:(?!</a:p>).)*?' + re.escape(_xml_escape(bullet4_ph)) + r'</a:t>(?:(?!</a:p>).)*?</a:p>)',
        flags=re.DOTALL,
    )
    m = pattern.search(xml)
    if not m:
        raise ValueError('Could not locate bullet-4 paragraph for extra_bullet injection.')
    bullet4_block = m.group(1)
    new_block = bullet4_block.replace(_xml_escape(bullet4_ph), _xml_escape(text))
    return xml.replace(bullet4_block, bullet4_block + new_block)


# =============================================================================
# Notes formatters
# =============================================================================

def format_concept_notes(video_script, key_terms=None, think_about=None, source_url=None):
    """Compose a notes string for concept/content slides.

    key_terms:   list of (term, definition) tuples, or None
    think_about: list of question strings (2 to 3), or None
    source_url:  authoritative source URL (NIST SP, MITRE, OWASP, RFC, etc.) or None
    """
    parts = ['VIDEO SCRIPT', '', video_script.strip()]
    if key_terms:
        parts += ['', 'KEY TERMS']
        parts += [f'{t}: {d}' for t, d in key_terms]
    if think_about:
        parts += ['', 'THINK ABOUT THIS']
        parts += [f'{i+1}. {q}' for i, q in enumerate(think_about)]
    if source_url:
        parts += ['', f'SOURCE: {source_url}']
    return '\n'.join(parts)


def format_title_notes(deck_id, deck_title, opening_line=None):
    """Compose title-slide notes: a brief cold-open guide."""
    default_opening = (
        f'Welcome to {deck_id}. Brief greeting, name the deck, then advance.'
    )
    return (
        f'TITLE SLIDE: {deck_id} \u2014 {deck_title}\n\n'
        f'This is the cold open. Brief greeting, name the deck, then immediately advance.\n\n'
        f'Suggested opening line:\n'
        f'{opening_line or default_opening}'
    )


def format_demo_notes(code, instructor_notes):
    """Compose demo-slide notes: the code (re-shown in notes) plus directing notes."""
    return f'DEMO COMMANDS (type or run live):\n\n{code.strip()}\n\n\nINSTRUCTOR NOTES\n\n{instructor_notes.strip()}'


def format_output_notes(output_text, instructor_notes):
    """Compose output-slide notes."""
    return f'EXPECTED OUTPUT\n\n{output_text.strip()}\n\n\nINSTRUCTOR NOTES\n\n{instructor_notes.strip()}'


# =============================================================================
# DeckBuilder
# =============================================================================

class DeckBuilder:
    """Assembles a slide deck from per-layout add_*_slide calls.

    Lifecycle:
        deck = DeckBuilder(template=..., png_out=..., work=...)
        deck.add_<layout>_slide(...)   # repeat for each slide
        deck.save(output_path)
    """

    def __init__(self, template=None, png_out=None, work=None):
        self.template = Path(template or DEFAULT_TEMPLATE)
        if not self.template.exists():
            raise FileNotFoundError(f'Template not found: {self.template}')

        self.png_out = Path(png_out or '/home/claude/build_pngs')
        self.png_out.mkdir(parents=True, exist_ok=True)

        self.work = Path(work or '/tmp/build_deck')
        if self.work.exists():
            shutil.rmtree(self.work)
        self.work.mkdir(parents=True)

        # Unpack the template
        self._run([sys.executable, str(UNPACK), str(self.template), str(self.work / 'unpacked')])
        self.unpacked = self.work / 'unpacked'
        self.slides_dir = self.unpacked / 'ppt' / 'slides'

        # Track new slides and pending PNG inserts / notes
        # Each slide entry: {'xml_file': 'slide36.xml', 'notes': '<text>', 'png': <png path or None>}
        self.slides = []

    # ----- subprocess helper -----
    def _run(self, cmd):
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            raise RuntimeError(f'Command failed: {cmd}\n{result.stderr}')
        return result.stdout

    # ----- internal: duplicate a template slide and return the new slide filename -----
    def _duplicate(self, layout_key):
        source_slide = f'slide{LAYOUT[layout_key]}.xml'
        out = self._run([sys.executable, str(ADD_SLIDE), str(self.unpacked), source_slide])
        m = re.search(r'Created (slide\d+\.xml)', out)
        if not m:
            raise RuntimeError(f'Could not parse add_slide output:\n{out}')
        return m.group(1)

    def _read(self, slide_filename):
        return (self.slides_dir / slide_filename).read_text()

    def _write(self, slide_filename, xml):
        (self.slides_dir / slide_filename).write_text(xml)

    # =========================================================================
    # Slide builders (one per layout)
    # =========================================================================

    def add_title_slide(self, course_id, course_name, subtitle, attribution=None, notes=''):
        """Title slide. course_id (large block) + course_name + subtitle + attribution line."""
        if attribution is None:
            attribution = PH['title']['attribution']
        fn = self._duplicate('title')
        xml = self._read(fn)
        # Note: course_id placeholder is 'IS3513'; if course_id matches, the replacement
        # is a no-op (harmless). We still call _replace_text so the surgical pattern is uniform.
        xml = _replace_text(xml, PH['title']['course_id'],   course_id)
        xml = _replace_text(xml, PH['title']['course_name'], course_name)
        xml = _replace_text(xml, PH['title']['subtitle'],    subtitle)
        xml = _replace_text(xml, PH['title']['attribution'], attribution)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_concept_slide(self, kicker, title, subhead, section_kicker,
                          card_heading, lead, bullets, extra_bullet=None, notes=''):
        """Concept slide: header keyword + title + yellow subhead + card with
        kicker/heading/lead and 4 bullets.

        Standard concept slides take exactly 4 bullets. The layout is sized for 4.

        extra_bullet (str, optional): exception parameter for the rare slide that
        genuinely needs a 5th item. When provided, a 5th paragraph is injected
        into the bullet text-body after the standard 4. Use this sparingly and
        only when none of the 5 items can be cut without losing meaning. The
        per-deck script should comment on why the exception applies.
        """
        if len(bullets) > 4:
            raise ValueError(f'Concept slide accepts up to 4 standard bullets, got {len(bullets)}. '
                             f'For a 5th item, use extra_bullet=...')
        bullets = list(bullets) + [''] * (4 - len(bullets))

        fn = self._duplicate('concept')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['concept']['kicker'],         kicker)
        xml = _replace_text(xml, PH['concept']['title'],          title)
        xml = _replace_text(xml, PH['concept']['subhead'],        subhead)
        xml = _replace_text(xml, PH['concept']['section_kicker'], section_kicker)
        xml = _replace_text(xml, PH['concept']['card_heading'],   card_heading)
        xml = _replace_text(xml, PH['concept']['lead'],           lead)
        # Inject the 5th bullet first, while the bullet-4 placeholder anchor is
        # still present in the XML. After this, the standard 4-bullet replacement
        # below will overwrite the 4th placeholder text and leave the new 5th
        # paragraph (which contains its own non-placeholder text) untouched.
        if extra_bullet:
            xml = _inject_extra_bullet(xml, extra_bullet)
        for ph, bullet in zip(PH['concept']['bullets'], bullets):
            xml = _replace_text(xml, ph, bullet)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_two_column_slide(self, kicker, title, subhead,
                              left_kicker, left_heading, left_bullets,
                              right_kicker, right_heading, right_bullets, notes=''):
        """Two-column comparison. Each side accepts up to 4 bullets."""
        if len(left_bullets) > 4 or len(right_bullets) > 4:
            raise ValueError('Two-column slide accepts up to 4 bullets per side.')
        left_bullets = list(left_bullets) + [''] * (4 - len(left_bullets))
        right_bullets = list(right_bullets) + [''] * (4 - len(right_bullets))

        fn = self._duplicate('two_col')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['two_col']['kicker'],        kicker)
        xml = _replace_text(xml, PH['two_col']['title'],         title)
        xml = _replace_text(xml, PH['two_col']['subhead'],       subhead)
        xml = _replace_text(xml, PH['two_col']['left_kicker'],   left_kicker)
        xml = _replace_text(xml, PH['two_col']['left_heading'],  left_heading)
        for ph, b in zip(PH['two_col']['left_bullets'], left_bullets):
            xml = _replace_text(xml, ph, b)
        xml = _replace_text(xml, PH['two_col']['right_kicker'],  right_kicker)
        xml = _replace_text(xml, PH['two_col']['right_heading'], right_heading)
        for ph, b in zip(PH['two_col']['right_bullets'], right_bullets):
            xml = _replace_text(xml, ph, b)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_intro_slide(self, kicker, title, subhead, role_kicker, name,
                        tagline, facts, pullquote, notes=''):
        """Intro slide: portrait area + role/name/tagline + 4 facts + pull-quote."""
        if len(facts) > 4:
            raise ValueError(f'Intro slide accepts up to 4 facts, got {len(facts)}.')
        facts = list(facts) + [''] * (4 - len(facts))

        fn = self._duplicate('intro')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['intro']['kicker'],      kicker)
        xml = _replace_text(xml, PH['intro']['title'],       title)
        xml = _replace_text(xml, PH['intro']['subhead'],     subhead)
        xml = _replace_text(xml, PH['intro']['role_kicker'], role_kicker)
        xml = _replace_text(xml, PH['intro']['name'],        name)
        xml = _replace_text(xml, PH['intro']['tagline'],     tagline)
        for ph, f in zip(PH['intro']['facts'], facts):
            xml = _replace_text(xml, ph, f)
        # Pull-quote uses entity-encoded smart quotes in v6 XML.
        xml = _replace_text_raw(xml, PH['intro']['pullquote_raw'], pullquote)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_team_slide(self, kicker, title, subhead, members, framing, notes=''):
        """Team slide: 5 mentor tiles + framing line.
        members: list of exactly 5 (name, role) tuples."""
        if len(members) != 5:
            raise ValueError(f'Team slide requires exactly 5 members, got {len(members)}.')

        fn = self._duplicate('team')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['team']['kicker'],  kicker)
        xml = _replace_text(xml, PH['team']['title'],   title)
        xml = _replace_text(xml, PH['team']['subhead'], subhead)
        # The 'First Lastname' and 'Role / Title' placeholders repeat 5 times.
        # Replace them sequentially using _replace_first so each member maps to one tile.
        for name, role in members:
            xml = _replace_first(xml, 'First Lastname', name)
            xml = _replace_first(xml, 'Role / Title',   role)
        xml = _replace_text(xml, PH['team']['framing'], framing)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_stats_slide(self, kicker, title, subhead, stats, lead, bullets, notes=''):
        """Stats slide: 3 stat tiles + supporting lead + up to 3 bullets.
        stats: list of exactly 3 (label, value) tuples."""
        if len(stats) != 3:
            raise ValueError(f'Stats slide requires exactly 3 stats, got {len(stats)}.')
        if len(bullets) > 3:
            raise ValueError(f'Stats slide accepts up to 3 bullets, got {len(bullets)}.')
        bullets = list(bullets) + [''] * (3 - len(bullets))

        fn = self._duplicate('stats')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['stats']['kicker'],  kicker)
        xml = _replace_text(xml, PH['stats']['title'],   title)
        xml = _replace_text(xml, PH['stats']['subhead'], subhead)
        # Replace each placeholder stat label/value with the supplied stat.
        for (old_label, old_value), (new_label, new_value) in zip(PH['stats']['stats'], stats):
            xml = _replace_text(xml, old_label, new_label)
            xml = _replace_text(xml, old_value, new_value)
        xml = _replace_text(xml, PH['stats']['lead'], lead)
        for ph, b in zip(PH['stats']['bullets'], bullets):
            xml = _replace_text(xml, ph, b)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_structure_slide(self, kicker, title, subhead, lead, cards, notes=''):
        """Three-card row. cards: list of exactly 3 (label, heading, body) tuples."""
        if len(cards) != 3:
            raise ValueError(f'Structure slide requires exactly 3 cards, got {len(cards)}.')

        fn = self._duplicate('structure')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['structure']['kicker'],  kicker)
        xml = _replace_text(xml, PH['structure']['title'],   title)
        xml = _replace_text(xml, PH['structure']['subhead'], subhead)
        xml = _replace_text(xml, PH['structure']['lead'],    lead)
        for (old_label, old_heading, old_body), (new_label, new_heading, new_body) in \
                zip(PH['structure']['cards'], cards):
            xml = _replace_text(xml, old_label,   new_label)
            xml = _replace_text(xml, old_heading, new_heading)
            xml = _replace_text(xml, old_body,    new_body)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_warning_slide(self, kicker, title, subhead, banner, banner_line,
                          rule_kicker, rule_oneliner, rule_body, bullets, notes=''):
        """Warning / hard-truth slide: banner + rule + up to 3 supporting bullets."""
        if len(bullets) > 3:
            raise ValueError(f'Warning slide accepts up to 3 bullets, got {len(bullets)}.')
        bullets = list(bullets) + [''] * (3 - len(bullets))

        fn = self._duplicate('warning')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['warning']['kicker'],         kicker)
        xml = _replace_text(xml, PH['warning']['title'],          title)
        xml = _replace_text(xml, PH['warning']['subhead'],        subhead)
        xml = _replace_text(xml, PH['warning']['banner'],         banner)
        xml = _replace_text(xml, PH['warning']['banner_line'],    banner_line)
        xml = _replace_text(xml, PH['warning']['rule_kicker'],    rule_kicker)
        xml = _replace_text(xml, PH['warning']['rule_oneliner'],  rule_oneliner)
        xml = _replace_text(xml, PH['warning']['rule_body'],      rule_body)
        for ph, b in zip(PH['warning']['bullets'], bullets):
            xml = _replace_text(xml, ph, b)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_support_slide(self, kicker, title, subhead,
                          philo_kicker, philo_heading, philo_body,
                          channels, notes=''):
        """Support slide: philosophy block + 4 channel tiles.
        channels: list of exactly 4 (label, value) tuples."""
        if len(channels) != 4:
            raise ValueError(f'Support slide requires exactly 4 channels, got {len(channels)}.')

        fn = self._duplicate('support')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['support']['kicker'],        kicker)
        xml = _replace_text(xml, PH['support']['title'],         title)
        xml = _replace_text(xml, PH['support']['subhead'],       subhead)
        xml = _replace_text(xml, PH['support']['philo_kicker'],  philo_kicker)
        xml = _replace_text(xml, PH['support']['philo_heading'], philo_heading)
        xml = _replace_text(xml, PH['support']['philo_body'],    philo_body)
        for (old_label, old_value), (new_label, new_value) in zip(PH['support']['channels'], channels):
            xml = _replace_text(xml, old_label, new_label)
            xml = _replace_text(xml, old_value, new_value)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_module_slide(self, number, title, labs_value, chap_value, client_value,
                          overview, labs, notes=''):
        """Module overview slide: big number + 3 stat tiles + overview + 3 labs.
        labs: list of exactly 3 (id, title, billing_type) tuples."""
        if len(labs) != 3:
            raise ValueError(f'Module slide requires exactly 3 labs, got {len(labs)}.')

        fn = self._duplicate('module')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['module']['number'],       str(number))
        xml = _replace_text(xml, PH['module']['title'],        title)
        xml = _replace_text(xml, PH['module']['labs_value'],   labs_value)
        xml = _replace_text(xml, PH['module']['chap_value'],   chap_value)
        xml = _replace_text(xml, PH['module']['client_value'], client_value)
        xml = _replace_text(xml, PH['module']['overview'],     overview)
        for (old_id, old_title, old_type), (new_id, new_title, new_type) in \
                zip(PH['module']['labs'], labs):
            xml = _replace_text(xml, old_id,    new_id)
            # Lab title placeholder repeats for Labs X.1 and X.2; use _replace_first.
            xml = _replace_first(xml, old_title, new_title)
            xml = _replace_first(xml, old_type,  new_type)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_guide_slide(self, kicker, title, subhead, callouts, notes=''):
        """Three stacked callouts. callouts: list of exactly 3 (label, heading, body) tuples."""
        if len(callouts) != 3:
            raise ValueError(f'Guide slide requires exactly 3 callouts, got {len(callouts)}.')

        fn = self._duplicate('guide')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['guide']['kicker'],  kicker)
        xml = _replace_text(xml, PH['guide']['title'],   title)
        xml = _replace_text(xml, PH['guide']['subhead'], subhead)
        for (old_label, old_heading, old_body), (new_label, new_heading, new_body) in \
                zip(PH['guide']['callouts'], callouts):
            xml = _replace_text(xml, old_label,   new_label)
            xml = _replace_text(xml, old_heading, new_heading)
            xml = _replace_text(xml, old_body,    new_body)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    def add_policy_slide(self, kicker, title, subhead, rules_kicker, rules, notes=''):
        """Four numbered rules. rules: list of exactly 4 (title, body) tuples.
        Rule numbers ('1'..'4') stay as-is in the template."""
        if len(rules) != 4:
            raise ValueError(f'Policy slide requires exactly 4 rules, got {len(rules)}.')

        fn = self._duplicate('policy')
        xml = self._read(fn)
        xml = _replace_text(xml, PH['policy']['kicker'],       kicker)
        xml = _replace_text(xml, PH['policy']['title'],        title)
        xml = _replace_text(xml, PH['policy']['subhead'],      subhead)
        xml = _replace_text(xml, PH['policy']['rules_kicker'], rules_kicker)
        # Title and body strings repeat 4 times; replace in document order.
        for rule_title, rule_body in rules:
            xml = _replace_first(xml, PH['policy']['rules_title_repeat'], rule_title)
            xml = _replace_first(xml, PH['policy']['rules_body_repeat'],  rule_body)
        self._write(fn, xml)
        self.slides.append({'xml_file': fn, 'notes': notes, 'png': None, 'is_png_slide': False})

    # =========================================================================
    # save: finalize sldIdLst, clean, pack, then inject notes
    # =========================================================================

    def save(self, output_path):
        """Update sldIdLst, clean, pack, then add speaker notes (and PNGs if any)."""
        # 1) Update sldIdLst to contain only our new slides in order
        pres_xml_path = self.unpacked / 'ppt' / 'presentation.xml'
        pres_xml = pres_xml_path.read_text()

        rels_xml = (self.unpacked / 'ppt' / '_rels' / 'presentation.xml.rels').read_text()
        slide_to_rid = {}
        for m in re.finditer(r'Id="(rId\d+)"\s+Type="[^"]*slide"\s+Target="slides/(slide\d+\.xml)"', rels_xml):
            slide_to_rid[m.group(2)] = m.group(1)

        entries = []
        for i, slide in enumerate(self.slides, start=400):
            rid = slide_to_rid.get(slide['xml_file'])
            if rid is None:
                raise RuntimeError(f"Could not find rId for {slide['xml_file']}")
            entries.append(f'    <p:sldId id="{i}" r:id="{rid}"/>')
        new_sldidlst = '<p:sldIdLst>\n' + '\n'.join(entries) + '\n  </p:sldIdLst>'

        pres_xml = re.sub(r'<p:sldIdLst>.*?</p:sldIdLst>', new_sldidlst, pres_xml, flags=re.DOTALL)
        pres_xml_path.write_text(pres_xml)

        # 2) Clean + pack to an intermediate file
        self._run([sys.executable, str(CLEAN), str(self.unpacked)])

        intermediate = self.work / '_intermediate.pptx'
        self._run([sys.executable, str(PACK), str(self.unpacked), str(intermediate),
                   '--original', str(self.template)])

        # 3) Open with python-pptx; add PNGs (if any future PNG slide types) and notes
        prs = Presentation(str(intermediate))
        if len(prs.slides) != len(self.slides):
            raise RuntimeError(
                f'Slide count mismatch: packed deck has {len(prs.slides)}, expected {len(self.slides)}'
            )

        for slide, spec in zip(prs.slides, self.slides):
            if spec['is_png_slide'] and spec['png']:
                slide.shapes.add_picture(spec['png'], PNG_X, PNG_Y, width=PNG_W)
            if spec['notes']:
                slide.notes_slide.notes_text_frame.text = spec['notes']

        # 4) Save final
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        print(f'Saved: {output_path}')
        print(f'  {len(self.slides)} slides, '
              f'{sum(1 for s in self.slides if s["png"])} PNGs, '
              f'{sum(1 for s in self.slides if s["notes"])} notes pages.')
